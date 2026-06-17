"""Generate filled proposal PowerPoint."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PURPLE    = RGBColor(0x5B, 0x21, 0xB6)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK      = RGBColor(0x1F, 0x17, 0x3A)
GRAY      = RGBColor(0x9C, 0xA3, 0xAF)
ORANGE    = RGBColor(0xF9, 0x73, 0x16)
LIGHT_BG  = RGBColor(0xF5, 0xF3, 0xFF)
TABLE_ROW = RGBColor(0xF3, 0xF4, 0xF6)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]


def rect(slide, x, y, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid() if fill else s.fill.background()
    if fill:
        s.fill.fore_color.rgb = fill
    s.line.fill.background() if not line else None
    if line:
        s.line.color.rgb = line
    return s


def tb(slide, text, x, y, w, h, size=11, bold=False, color=DARK,
       align=PP_ALIGN.LEFT, italic=False):
    t = slide.shapes.add_textbox(x, y, w, h)
    tf = t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.color.rgb = color
    r.font.name  = "Sarabun"
    r.font.italic = italic
    return t


def field(slide, label, content, x, y, w, h, placeholder=False):
    """Labeled field box. placeholder=True → gray italic text."""
    rect(slide, x, y, w, h, fill=WHITE, line=RGBColor(0xD1, 0xD5, 0xDB))
    if label:
        tb(slide, label, x + Inches(0.1), y + Inches(0.04),
           w - Inches(0.2), Inches(0.22), size=9, bold=True)
    content_y = y + (Inches(0.26) if label else Inches(0.08))
    tb(slide, content, x + Inches(0.1), content_y,
       w - Inches(0.2), h - Inches(0.32 if label else 0.15),
       size=10, color=GRAY if placeholder else DARK, italic=placeholder)


def section(slide, num, title, x, y, w):
    tb(slide, f"{num}. {title}", x, y, w, Inches(0.3), size=13, bold=True, color=PURPLE)
    rect(slide, x, y + Inches(0.29), w, Pt(2), fill=PURPLE)


def header(slide):
    rect(slide, 0, 0, W, Inches(0.07), fill=PURPLE)
    rect(slide, 0, Inches(0.07), Inches(0.12), Inches(1.0), fill=PURPLE)


def slide_top(slide, subtitle=""):
    header(slide)
    tb(slide, "Project Pitching Template | วิชาที่ 7",
       Inches(0.2), Inches(0.12), Inches(10), Inches(0.45), size=22, bold=True)
    tb(slide, "AI Engineering / LLM / RAG / AI Workflow" + (f"  •  {subtitle}" if subtitle else ""),
       Inches(0.2), Inches(0.55), Inches(10), Inches(0.25), size=10, color=PURPLE)
    rect(slide, Inches(0.2), Inches(0.78), W - Inches(0.4), Pt(1),
         fill=RGBColor(0xE5, 0xE7, 0xEB))


def purple_table_header(slide, headers, x, y, w, col_w=None):
    n = len(headers)
    cw = col_w or [w / n] * n
    cx = x
    h = Inches(0.35)
    for i, hdr in enumerate(headers):
        rect(slide, cx, y, cw[i], h, fill=PURPLE)
        tb(slide, hdr, cx + Inches(0.08), y + Inches(0.05),
           cw[i] - Inches(0.1), h, size=9, bold=True, color=WHITE)
        cx += cw[i]
    return y + h


# ═══════════════════════════════════════════════════════════
# SLIDE 1 — Submission
# ═══════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
slide_top(s1)
tb(s1, "แบบฟอร์มนำเสนอโปรเจคต์ — กรอกข้อมูลให้ครบก่อนส่ง",
   Inches(0.2), Inches(0.85), Inches(12), Inches(0.3), size=11, bold=True)
rect(s1, Inches(0.5), Inches(1.25), Inches(12.33), Inches(4.0), fill=LIGHT_BG)

FW = Inches(5.5)
FH = Inches(1.1)

field(s1, "Project Topic (หัวข้อ):", "Project Knowledge Hub",
      Inches(0.7), Inches(1.45), FW, FH)
field(s1, "Topic Number:", "",
      Inches(6.9), Inches(1.45), FW, FH, placeholder=True)
field(s1, "Group Name:", "",
      Inches(0.7), Inches(2.7), FW, FH, placeholder=True)
field(s1, "Date Submitted:", "17/06/2026",
      Inches(6.9), Inches(2.7), FW, FH)

# ═══════════════════════════════════════════════════════════
# SLIDE 2 — Overview, Problem, Objectives
# ═══════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
slide_top(s2, "Overview, problem statement, objectives")

LX = Inches(0.2)
RX = Inches(6.85)
CW = Inches(6.4)

section(s2, 1, "Project Overview (ภาพรวมโปรเจคต์)", LX, Inches(0.9), CW)
field(s2, "Project Title:", "Project Knowledge Hub",
      LX, Inches(1.25), CW, Inches(0.65))
field(s2, "Project Description:",
      "RAG-based Q&A assistant ที่ตอบคำถามจากเอกสารโปรเจคต์ "
      "พร้อม citation อ้างอิงแหล่งที่มา เมื่อไม่พบคำตอบในเอกสาร "
      "ระบบจะ abstain แทนการเดา",
      LX, Inches(1.95), CW, Inches(0.85))
field(s2, "Business Context:",
      "ทีมวิศวกรรมซอฟต์แวร์ที่มีเอกสาร (spec, runbook, ADR, meeting notes) "
      "กระจายอยู่หลายที่",
      LX, Inches(2.85), CW, Inches(0.75))

section(s2, 2, "Problem Statement (ปัญหาที่ต้องแก้ไข)", LX, Inches(3.72), CW)
field(s2, "Current Problem:",
      "ทีม Engineering ค้นหาข้อมูลจากเอกสารหลายไฟล์ด้วยตนเอง "
      "ไม่มีระบบกลางในการตอบคำถาม",
      LX, Inches(4.07), CW, Inches(0.75))
field(s2, "Impact of Problem:",
      "เสียเวลาค้นหา, ตัดสินใจโดยข้อมูลไม่ครบ, เกิด knowledge gap ในทีม",
      LX, Inches(4.87), CW, Inches(0.65))
field(s2, "Proposed Solution:",
      "RAG-based Q&A ที่ตอบจากเอกสารจริง พร้อม citation "
      "และไม่ hallucinate เมื่อข้อมูลไม่เพียงพอ",
      LX, Inches(5.57), CW, Inches(0.8))

section(s2, 3, "Objectives (วัตถุประสงค์) — SMART Goals", RX, Inches(0.9), CW)
tb(s2, "ระบุวัตถุประสงค์ 3-5 ข้อ (SMART Goals)",
   RX, Inches(1.24), CW, Inches(0.22), size=9, color=GRAY, italic=True)

objectives = [
    "ระบบ ingest เอกสาร .md และ .txt และ query ผลลัพธ์ได้ภายใน 3 วินาที",
    "คำตอบทุกข้อระบุ source_file และ section_heading เป็น citation",
    "ระบบ abstain (ตอบว่า 'ไม่พบข้อมูล') เมื่อ similarity score ต่ำกว่า threshold",
    "รองรับเอกสารทั้งภาษาไทยและอังกฤษผ่าน multilingual-e5-small embedding",
    "Deploy บน Hugging Face Spaces เข้าถึงได้ผ่าน browser โดยไม่ติดตั้งเพิ่ม",
]
for i, obj in enumerate(objectives):
    y = Inches(1.5) + i * Inches(0.82)
    field(s2, f"{i+1}.", obj, RX, y, CW, Inches(0.75))

# Tip
tip_box = rect(s2, RX, Inches(5.72), CW, Inches(0.68),
               fill=RGBColor(0xFF, 0xF7, 0xED), line=ORANGE)
tb(s2, "Tip", RX + Inches(0.1), Inches(5.79),
   Inches(0.4), Inches(0.25), size=9, bold=True, color=ORANGE)
tb(s2, "Use measurable outcomes where possible, e.g., reduce SLA breach detection time or improve report accuracy.",
   RX + Inches(0.5), Inches(5.79), CW - Inches(0.6), Inches(0.6),
   size=9, color=RGBColor(0x92, 0x40, 0x08))

# ═══════════════════════════════════════════════════════════
# SLIDE 3 — Scope, AI/LLM/RAG, Workflow header
# ═══════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
slide_top(s3, "Scope, AI / LLM / RAG Components, workflow")

section(s3, 4, "Scope (ขอบเขตโปรเจคต์)", Inches(0.2), Inches(0.9), Inches(12.9))

y_scope = purple_table_header(s3, ["In Scope", "Out of Scope"],
                               Inches(0.2), Inches(1.25), Inches(12.9))
in_scope = [
    "Ingest เอกสาร .md และ .txt",
    "RAG query พร้อม citation (source_file, section_heading)",
    "Streamlit chat UI, deploy บน Hugging Face Spaces",
]
out_scope = [
    "Fine-tuning LLM",
    "รองรับไฟล์ PDF, DOCX หรือรูปภาพ",
    "User authentication / multi-tenant",
]
HW = Inches(12.9) / 2
for i, (ins, out) in enumerate(zip(in_scope, out_scope)):
    bg = TABLE_ROW if i % 2 == 0 else WHITE
    rect(s3, Inches(0.2),     y_scope + i*Inches(0.5), HW,        Inches(0.5), fill=bg)
    rect(s3, Inches(0.2)+HW,  y_scope + i*Inches(0.5), HW,        Inches(0.5), fill=bg)
    tb(s3, ins,  Inches(0.35),     y_scope + i*Inches(0.5)+Inches(0.08), HW-Inches(0.3), Inches(0.4), size=10)
    tb(s3, out,  Inches(0.35)+HW,  y_scope + i*Inches(0.5)+Inches(0.08), HW-Inches(0.3), Inches(0.4), size=10)

section(s3, 5, "AI / LLM / RAG Components (องค์ประกอบของระบบ AI)",
        Inches(0.2), Inches(3.7), Inches(12.9))

cw3 = Inches(12.9) / 3
y5 = purple_table_header(s3, ["LLM Model Used", "RAG / Knowledge Base", "Agent / Workflow"],
                          Inches(0.2), Inches(4.05), Inches(12.9))
row5 = ["GPT-4o-mini (GitHub Models,\nOpenAI-compatible endpoint)",
        "เอกสาร .md/.txt + ChromaDB\n+ intfloat/multilingual-e5-small",
        "Single-turn Q&A\n(no agentic loop)"]
for i, val in enumerate(row5):
    rect(s3, Inches(0.2)+i*cw3, y5, cw3, Inches(0.75), fill=TABLE_ROW if i%2==0 else WHITE)
    tb(s3, val, Inches(0.3)+i*cw3, y5+Inches(0.08), cw3-Inches(0.2), Inches(0.65), size=10)

cw4 = Inches(12.9) / 4
y4h = y5 + Inches(0.75) + Inches(0.12)
y4r = purple_table_header(s3, ["Input", "AI Processing", "Output", "Human Review"],
                           Inches(0.2), y4h, Inches(12.9))
row4 = [
    "คำถามจากผู้ใช้\nผ่าน Streamlit UI",
    "Embed คำถาม → ค้นหา top-k chunks\nจาก ChromaDB → ส่ง context+คำถามให้ LLM",
    "คำตอบ + citation\n(source_file, section_heading)",
    "ไม่มี human-in-the-loop\nระบบ abstain อัตโนมัติ",
]
for i, val in enumerate(row4):
    rect(s3, Inches(0.2)+i*cw4, y4r, cw4, Inches(0.75), fill=TABLE_ROW if i%2==0 else WHITE)
    tb(s3, val, Inches(0.3)+i*cw4, y4r+Inches(0.06), cw4-Inches(0.2), Inches(0.68), size=9)

y_risk = y4r + Inches(0.75) + Inches(0.05)
field(s3, "Hallucination / Risk Mitigation:",
      "ใช้ similarity score threshold — ถ้าต่ำกว่า threshold ตอบ 'ไม่พบข้อมูลในเอกสาร' "
      "แทน; system prompt สั่ง LLM ให้ตอบเฉพาะจาก context ที่ให้เท่านั้น",
      Inches(0.2), y_risk, Inches(12.9), Inches(0.65))

section(s3, 6, "Workflow / Architecture (ขั้นตอนการทำงาน)",
        Inches(0.2), y_risk + Inches(0.7), Inches(12.9))

# ═══════════════════════════════════════════════════════════
# SLIDE 4 — Workflow steps + Tools + Outcomes + Timeline + Team
# ═══════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
slide_top(s4, "Tools, outcomes, timeline, team")

tb(s4, "อธิบาย flow ของระบบตั้งแต่ต้นจนจบ (Input → Process → Output)",
   Inches(0.2), Inches(0.9), Inches(12.9), Inches(0.22), size=9, color=GRAY, italic=True)

steps = [
    "User อัปโหลด\nเอกสาร .md/.txt",
    "Parse + chunk\ntext",
    "Embed chunks\n(multilingual-e5-small)",
    "Store embeddings\nใน ChromaDB",
    "User ถามคำถาม\n→ embed → top-k",
    "LLM สร้างคำตอบ\n+ citation",
]
sw = Inches(12.9) / 6
ys = Inches(1.15)
y_step_r = purple_table_header(s4, [f"Step {i+1}" for i in range(6)],
                                Inches(0.2), ys, Inches(12.9))
for i, step in enumerate(steps):
    rect(s4, Inches(0.2)+i*sw, y_step_r, sw, Inches(0.75),
         fill=TABLE_ROW if i%2==0 else WHITE)
    tb(s4, step, Inches(0.28)+i*sw, y_step_r+Inches(0.06),
       sw-Inches(0.15), Inches(0.68), size=9)

field(s4, "Architecture Notes:",
      "Embedding model โหลดครั้งเดียว (singleton) เพื่อลด latency; "
      "ChromaDB persistent path อ่านจาก CHROMA_DB_PATH env var",
      Inches(0.2), y_step_r+Inches(0.78), Inches(12.9), Inches(0.6))

section(s4, 7, "Tools & Technologies", Inches(0.2), Inches(2.65), Inches(12.9))
cw7 = Inches(12.9) / 4
y7r = purple_table_header(s4,
    ["Programming\nLanguage", "Database / Storage", "Frameworks /\nLibraries", "Visualization /\nOutput"],
    Inches(0.2), Inches(3.0), Inches(12.9))
tools = [
    "Python 3.12",
    "ChromaDB\n(persistent vector store)",
    "Streamlit, sentence-transformers,\nopenai SDK, chromadb",
    "Streamlit web UI\n(Hugging Face Spaces)",
]
for i, t in enumerate(tools):
    rect(s4, Inches(0.2)+i*cw7, y7r, cw7, Inches(0.75), fill=TABLE_ROW if i%2==0 else WHITE)
    tb(s4, t, Inches(0.3)+i*cw7, y7r+Inches(0.08), cw7-Inches(0.2), Inches(0.65), size=10)

section(s4, 8, "Expected Outcomes (ผลลัพธ์ที่คาดหวัง)",
        Inches(0.2), Inches(4.05), Inches(12.9))
field(s4, "Deliverables:",
      "Streamlit web app บน Hugging Face Spaces, ingestion pipeline (.md/.txt), RAG Q&A พร้อม citation",
      Inches(0.2), Inches(4.4), Inches(12.9), Inches(0.6))
field(s4, "Success Metrics:",
      "Retrieval precision ≥ 80% (manual test set 20 Q&A pairs), "
      "abstention ถูกต้องเมื่อถามนอกเอกสาร, response time < 5 วินาที",
      Inches(0.2), Inches(5.05), Inches(12.9), Inches(0.6))

section(s4, 9, "Timeline (แผนการดำเนินงาน)", Inches(0.2), Inches(5.78), Inches(12.9))
cw9 = Inches(12.9) / 4
y9r = purple_table_header(s4,
    ["สัปดาห์ที่ 1-2", "สัปดาห์ที่ 3-4", "สัปดาห์ที่ 5-6", "สัปดาห์ที่ 7-8"],
    Inches(0.2), Inches(6.1), Inches(12.9))
timeline = [
    "RAG query pipeline\n+ wire กับ Streamlit UI",
    "Chat UI + citation display\n+ abstention logic",
    "Document upload UI\n+ end-to-end testing",
    "Evaluation (20 Q&A test set)\n+ final deployment",
]
for i, t in enumerate(timeline):
    rect(s4, Inches(0.2)+i*cw9, y9r, cw9, Inches(0.7), fill=TABLE_ROW if i%2==0 else WHITE)
    tb(s4, t, Inches(0.3)+i*cw9, y9r+Inches(0.06), cw9-Inches(0.2), Inches(0.62), size=9)

section(s4, 10, "Team Members (สมาชิกในกลุ่ม)", Inches(0.2), Inches(7.05), Inches(12.9))

# ═══════════════════════════════════════════════════════════
# SLIDE 5 — Team Members table (overflow)
# ═══════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank)
slide_top(s5, "Team Members")

col_w = [Inches(0.5), Inches(4.5), Inches(3.5), Inches(4.4)]
y_tm = purple_table_header(s5, ["#", "ชื่อ - นามสกุล", "รหัสนักศึกษา", "บทบาทใน Project"],
                            Inches(0.2), Inches(1.0), Inches(12.9), col_w=col_w)
members = [
    ("1", "นาย นราวิชญ์ สมนา",         "", "ผู้ให้ requirement และดูภาพรวม"),
    ("2", "นาย ภัทรงค์ สมบัติแก้ว",    "", "ผู้ให้คำแนะนำ"),
    ("3", "นาย อภิสิทธิ์ วงศ์วิศิษฐ์", "", "ผู้ออกแบบและพัฒนาระบบ"),
]
for row_i, (num, name, sid, role) in enumerate(members):
    bg = TABLE_ROW if row_i % 2 == 0 else WHITE
    cx = Inches(0.2)
    for ci, (val, cw) in enumerate(zip([num, name, sid, role], col_w)):
        rect(s5, cx, y_tm + row_i*Inches(0.6), cw, Inches(0.6), fill=bg)
        tb(s5, val, cx+Inches(0.1), y_tm + row_i*Inches(0.6)+Inches(0.1),
           cw-Inches(0.15), Inches(0.45), size=11)
        cx += cw

out = "/home/nate_river/project/AI-for-engineering-course/ai-assistant-project/proposal_template.pptx"
prs.save(out)
print(f"Saved: {out}")
